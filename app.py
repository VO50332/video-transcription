import os
import uuid
import threading
import tempfile
import json
from pathlib import Path
from flask import Flask, request, jsonify, render_template, send_file

app = Flask(__name__)

OUTPUTS_DIR = Path(__file__).parent / "outputs"
OUTPUTS_DIR.mkdir(exist_ok=True)

# In-memory job store: job_id -> {"status", "progress", "transcript", "summary", "pptx_path", "txt_path", "error"}
jobs = {}


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def extract_audio(video_url: str, cookies: str | None, out_path: str) -> None:
    """Use yt-dlp to extract audio from the video URL into an mp3 file."""
    import yt_dlp

    ydl_opts = {
        "format": "bestaudio/best",
        "outtmpl": out_path,
        "postprocessors": [{
            "key": "FFmpegExtractAudio",
            "preferredcodec": "mp3",
            "preferredquality": "96",
        }],
        "quiet": True,
        "no_warnings": True,
    }

    if cookies:
        # Write cookies to a temp file
        cookie_file = out_path + ".cookies.txt"
        with open(cookie_file, "w") as f:
            f.write(cookies)
        ydl_opts["cookiefile"] = cookie_file

    with yt_dlp.YoutubeDL(ydl_opts) as ydl:
        ydl.download([video_url])


def transcribe_audio(audio_path: str, model_size: str = "base") -> str:
    """Transcribe audio using OpenAI Whisper (runs locally)."""
    import whisper
    model = whisper.load_model(model_size)
    result = model.transcribe(audio_path, fp16=False)
    return result["text"].strip()


def generate_summary_and_slides(transcript: str, video_url: str) -> tuple[str, list[dict]]:
    """Use Claude API to produce a summary and slide content from the transcript."""
    import anthropic

    client = anthropic.Anthropic()

    prompt = f"""You are an expert content analyst. Below is the full transcript of an online video.

TRANSCRIPT:
{transcript[:40000]}

Your task:
1. Write a concise but thorough **summary** (3-5 paragraphs) covering the key points, insights, and conclusions.
2. Create a structured **PowerPoint outline** with 8-12 slides. For each slide provide:
   - A short title (max 8 words)
   - 3-5 bullet points (each max 15 words)
   - Optional speaker notes (1-2 sentences for context)

Return ONLY valid JSON in this exact schema (no markdown fences):
{{
  "summary": "<multi-paragraph summary as plain text>",
  "slides": [
    {{
      "title": "Slide Title",
      "bullets": ["bullet 1", "bullet 2", "bullet 3"],
      "notes": "Optional speaker notes."
    }}
  ]
}}"""

    message = client.messages.create(
        model="claude-opus-4-6",
        max_tokens=4096,
        messages=[{"role": "user", "content": prompt}],
    )

    raw = message.content[0].text.strip()
    # Strip accidental markdown fences
    if raw.startswith("```"):
        raw = raw.split("```", 2)[1]
        if raw.startswith("json"):
            raw = raw[4:]
        raw = raw.rsplit("```", 1)[0].strip()

    data = json.loads(raw)
    return data["summary"], data["slides"]


def build_pptx(slides: list[dict], summary: str, output_path: str, video_url: str) -> None:
    """Build a .pptx file from the slide data."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Theme colors
    BG_COLOR = RGBColor(0x1A, 0x1A, 0x2E)
    ACCENT = RGBColor(0xE9, 0x4F, 0x37)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)

    def set_slide_bg(slide, color: RGBColor):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_textbox(slide, text, left, top, width, height,
                    font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        return txBox

    # ---- Title slide ----
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide, BG_COLOR)

    # Accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(3.2), Inches(13.33), Inches(0.07))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    add_textbox(slide, "Video Transcript & Analysis",
                Inches(1), Inches(1.2), Inches(11), Inches(1.2),
                font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, video_url,
                Inches(1), Inches(3.5), Inches(11), Inches(0.6),
                font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_textbox(slide, "Auto-generated with Claude AI + Whisper",
                Inches(1), Inches(4.3), Inches(11), Inches(0.5),
                font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # ---- Summary slide ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)

    # Title bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    add_textbox(slide, "Summary", Inches(0.4), Inches(0.15), Inches(12), Inches(0.8),
                font_size=28, bold=True, color=WHITE)

    # Summary text (truncated if too long for slide)
    summary_short = summary[:1200] + ("…" if len(summary) > 1200 else "")
    add_textbox(slide, summary_short,
                Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8),
                font_size=14, color=LIGHT_GRAY, wrap=True)

    # ---- Content slides ----
    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide, BG_COLOR)

        # Title bar
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()

        add_textbox(slide, slide_data["title"],
                    Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.8),
                    font_size=28, bold=True, color=WHITE)

        # Bullets
        bullet_top = Inches(1.25)
        bullet_height = Inches(0.55)
        for i, bullet in enumerate(slide_data.get("bullets", [])):
            add_textbox(slide, f"• {bullet}",
                        Inches(0.7), bullet_top + i * bullet_height,
                        Inches(11.9), bullet_height + Inches(0.1),
                        font_size=17, color=LIGHT_GRAY)

        # Speaker notes
        notes_text = slide_data.get("notes", "")
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text

    # ---- Thank You slide ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(3.2), Inches(13.33), Inches(0.07))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    add_textbox(slide, "Thank You", Inches(1), Inches(2.2), Inches(11), Inches(1.2),
                font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Generated by Video Transcriber",
                Inches(1), Inches(3.5), Inches(11), Inches(0.6),
                font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    prs.save(output_path)


# ---------------------------------------------------------------------------
# Background job runner
# ---------------------------------------------------------------------------

def run_job(job_id: str, video_url: str, cookies: str | None, whisper_model: str):
    job = jobs[job_id]

    with tempfile.TemporaryDirectory() as tmpdir:
        audio_base = os.path.join(tmpdir, "audio")
        audio_file = audio_base + ".mp3"

        try:
            # Step 1: extract audio
            job["progress"] = "Extracting audio from video…"
            extract_audio(video_url, cookies, audio_base)

            # yt-dlp may append extension differently
            if not os.path.exists(audio_file):
                candidates = [f for f in os.listdir(tmpdir) if f.startswith("audio")]
                if candidates:
                    audio_file = os.path.join(tmpdir, candidates[0])
                else:
                    raise FileNotFoundError("Audio extraction produced no output file.")

            # Step 2: transcribe
            job["progress"] = f"Transcribing audio with Whisper ({whisper_model})… this may take a few minutes."
            transcript = transcribe_audio(audio_file, whisper_model)
            job["transcript"] = transcript

            # Step 3: Claude summary + slides
            job["progress"] = "Generating summary and slides with Claude AI…"
            summary, slides = generate_summary_and_slides(transcript, video_url)
            job["summary"] = summary

            # Step 4: build pptx
            job["progress"] = "Building PowerPoint presentation…"
            pptx_path = str(OUTPUTS_DIR / f"{job_id}.pptx")
            build_pptx(slides, summary, pptx_path, video_url)
            job["pptx_path"] = pptx_path

            # Step 5: save transcript txt
            txt_path = str(OUTPUTS_DIR / f"{job_id}_transcript.txt")
            with open(txt_path, "w") as f:
                f.write(f"VIDEO URL: {video_url}\n\n")
                f.write("=== SUMMARY ===\n\n")
                f.write(summary + "\n\n")
                f.write("=== FULL TRANSCRIPT ===\n\n")
                f.write(transcript)
            job["txt_path"] = txt_path

            job["status"] = "done"
            job["progress"] = "Complete!"

        except Exception as e:
            job["status"] = "error"
            job["error"] = str(e)
            job["progress"] = f"Error: {e}"


# ---------------------------------------------------------------------------
# Routes
# ---------------------------------------------------------------------------

@app.route("/")
def index():
    return render_template("index.html")


@app.route("/transcribe", methods=["POST"])
def transcribe():
    data = request.get_json()
    video_url = (data.get("url") or "").strip()
    cookies = (data.get("cookies") or "").strip() or None
    whisper_model = data.get("model", "base")

    if not video_url:
        return jsonify({"error": "No URL provided"}), 400

    job_id = str(uuid.uuid4())
    jobs[job_id] = {
        "status": "running",
        "progress": "Starting…",
        "transcript": None,
        "summary": None,
        "pptx_path": None,
        "txt_path": None,
        "error": None,
    }

    t = threading.Thread(target=run_job, args=(job_id, video_url, cookies, whisper_model), daemon=True)
    t.start()

    return jsonify({"job_id": job_id})


@app.route("/status/<job_id>")
def status(job_id):
    job = jobs.get(job_id)
    if not job:
        return jsonify({"error": "Unknown job"}), 404
    return jsonify({
        "status": job["status"],
        "progress": job["progress"],
        "has_transcript": job["transcript"] is not None,
        "has_summary": job["summary"] is not None,
        "error": job["error"],
    })


@app.route("/result/<job_id>")
def result(job_id):
    job = jobs.get(job_id)
    if not job or job["status"] != "done":
        return jsonify({"error": "Not ready"}), 404
    return jsonify({
        "transcript": job["transcript"],
        "summary": job["summary"],
    })


@app.route("/download/pptx/<job_id>")
def download_pptx(job_id):
    job = jobs.get(job_id)
    if not job or not job.get("pptx_path"):
        return jsonify({"error": "Not ready"}), 404
    return send_file(job["pptx_path"], as_attachment=True,
                     download_name="video_presentation.pptx",
                     mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")


@app.route("/download/txt/<job_id>")
def download_txt(job_id):
    job = jobs.get(job_id)
    if not job or not job.get("txt_path"):
        return jsonify({"error": "Not ready"}), 404
    return send_file(job["txt_path"], as_attachment=True,
                     download_name="transcript_and_summary.txt",
                     mimetype="text/plain")


if __name__ == "__main__":
    app.run(debug=False, host="0.0.0.0", port=5000)
